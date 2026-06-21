from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parents[1]
ASSET_DIR = ROOT / "00_Notes" / "real_app_screenshots"
OUT = ROOT / "BigSam_Steganography_Internal_Defense_REAL_Screenshots_Projector_Readable.pptx"

NAVY = RGBColor(15, 23, 42)
SLATE = RGBColor(51, 65, 85)
MUTED = RGBColor(100, 116, 139)
LIGHT = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)
BLUE = RGBColor(37, 99, 235)
CYAN = RGBColor(6, 182, 212)
TEAL = RGBColor(20, 184, 166)
GREEN = RGBColor(16, 185, 129)
PURPLE = RGBColor(124, 58, 237)
AMBER = RGBColor(245, 158, 11)
BORDER = RGBColor(203, 213, 225)
FONT_HEAD = "Aptos Display"
FONT_BODY = "Aptos"


def fill(shape, color, transparency=0):
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.fill.transparency = transparency

def line(shape, color=None, width=1):
    if color is None: shape.line.fill.background()
    else:
        shape.line.color.rgb = color; shape.line.width = Pt(width)

def textbox(slide, text, x, y, w, h, size=22, color=SLATE, bold=False, align=PP_ALIGN.LEFT, font=FONT_BODY, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True; tf.vertical_anchor = valign
    tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04); tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; r.font.name = font; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return box

def bg(slide, dark=False):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = NAVY if dark else LIGHT
    # soft corner accents
    for i, c in enumerate([CYAN, BLUE, TEAL]):
        shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.4 + i*.45), Inches(-.9 + i*.12), Inches(2.3), Inches(2.3))
        fill(shp, c, 74 if not dark else 82); line(shp)

def title(slide, title, subtitle=None, dark=False):
    color = WHITE if dark else NAVY
    textbox(slide, title, .65, .35, 11.7, .6, size=34, color=color, bold=True, font=FONT_HEAD)
    if subtitle:
        textbox(slide, subtitle, .68, .96, 10.8, .38, size=17, color=(RGBColor(203,213,225) if dark else MUTED))
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(.68), Inches(1.36), Inches(1.35), Inches(.08))
    fill(bar, CYAN if dark else BLUE); line(bar)

def footer(slide, n, total=12, dark=False):
    color = RGBColor(148,163,184) if dark else MUTED
    textbox(slide, "Secure Web-Based Medical Image Management System using Steganography", .65, 7.12, 8.8, .24, size=10, color=color)
    textbox(slide, f"{n:02d}/{total:02d}", 12.0, 7.08, .65, .25, size=10, color=color, align=PP_ALIGN.RIGHT)

def card(slide, x, y, w, h, heading, body, accent=BLUE, heading_size=22, body_size=18, fill_color=WHITE):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    fill(shp, fill_color); line(shp, BORDER, 1.0)
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(.13), Inches(h))
    fill(strip, accent); line(strip)
    textbox(slide, heading, x+.33, y+.20, w-.55, .40, size=heading_size, color=NAVY, bold=True)
    textbox(slide, body, x+.33, y+.75, w-.55, h-.92, size=body_size, color=SLATE)

def bullets(slide, items, x, y, w, size=23, color=SLATE, accent=BLUE, gap=.62):
    for i, item in enumerate(items):
        yy = y + i*gap
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(yy+.13), Inches(.14), Inches(.14))
        fill(dot, accent); line(dot)
        textbox(slide, item, x+.32, yy, w-.32, .46, size=size, color=color)

def pill(slide, text, x, y, w, color=BLUE):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(.42))
    fill(shp, color); line(shp)
    textbox(slide, text, x+.05, y+.095, w-.1, .20, size=12.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

def add_picture_fit(slide, path, x, y, w, h):
    # Fit within bounding box while preserving aspect ratio
    from PIL import Image
    im = Image.open(path)
    iw, ih = im.size
    box_ratio = w/h; img_ratio = iw/ih
    if img_ratio >= box_ratio:
        pic_w = w; pic_h = w / img_ratio
    else:
        pic_h = h; pic_w = h * img_ratio
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(pic_w), height=Inches(pic_h))

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]
TOTAL = 12

# 1 Title
s = prs.slides.add_slide(blank); bg(s, True)
textbox(s, "INTERNAL DEFENSE PRESENTATION", .72, .72, 5.2, .3, size=14, color=CYAN, bold=True)
textbox(s, "Development of a Secure Web-Based Medical Image Management System Using Steganography", .72, 1.35, 11.0, 1.55, size=42, color=WHITE, bold=True, font=FONT_HEAD)
textbox(s, "Protecting patient diagnosis data inside medical images without visible distortion", .76, 3.05, 10.2, .45, size=22, color=RGBColor(203,213,225))
card(s, .78, 5.05, 5.7, 1.08, "Presenter", "OMONIYI SAMSON DAMILARE\nDepartment of Computer Science", CYAN, heading_size=16, body_size=18, fill_color=RGBColor(30,41,59))
# overlay readable light text for dark card
textbox(s, "Presenter", 1.08, 5.25, 2.0, .25, size=16, color=CYAN, bold=True)
textbox(s, "OMONIYI SAMSON DAMILARE\nDepartment of Computer Science", 1.08, 5.58, 4.8, .42, size=17, color=WHITE, bold=True)
footer(s, 1, TOTAL, True)

# 2 Outline
s = prs.slides.add_slide(blank); bg(s); title(s, "Presentation Outline", "Large-type version for dim projectors and back-row readability")
outline = ["Introduction", "Literature Review", "Motivation", "Objectives", "Methodology", "Results & Discussion", "Contribution", "Conclusion & Recommendations", "References"]
for i, item in enumerate(outline):
    x = .9 + (i % 3)*4.05; y = 1.78 + (i//3)*1.38
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.45), Inches(.88))
    fill(shp, WHITE); line(shp, BORDER, 1)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+.18), Inches(y+.2), Inches(.48), Inches(.48)); fill(circ, BLUE if i < 5 else TEAL); line(circ)
    textbox(s, str(i+1), x+.305, y+.30, .22, .18, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    textbox(s, item, x+.83, y+.24, 2.4, .35, size=20, color=NAVY, bold=True)
footer(s, 2, TOTAL)

# 3 Introduction
s = prs.slides.add_slide(blank); bg(s); title(s, "Introduction", "Medical image security problem")
card(s, .8, 1.65, 3.65, 3.95, "Medical images carry context", "X-rays and scans are usually shared with notes, findings, and patient identifiers.", BLUE)
card(s, 4.85, 1.65, 3.65, 3.95, "Steganography hides presence", "Sensitive text is embedded inside the image itself, while the image still looks normal.", CYAN)
card(s, 8.9, 1.65, 3.65, 3.95, "This project", "A secure web prototype for embedding and recovering patient diagnosis data from medical images.", TEAL)
pill(s, "Confidentiality", 1.1, 6.1, 1.75, BLUE); pill(s, "Image Quality", 3.15, 6.1, 1.55, CYAN); pill(s, "Audit Logging", 5.0, 6.1, 1.55, TEAL); pill(s, "Web Access", 6.85, 6.1, 1.35, PURPLE)
footer(s, 3, TOTAL)

# 4 Literature Review
s = prs.slides.add_slide(blank); bg(s); title(s, "Review of Literature and Related Works", "Objective • Methodology • Result • Contribution")
card(s, .8, 1.65, 3.65, 4.45, "Image steganography studies", "Prior work shows the key trade-off: hiding capacity versus image imperceptibility and robustness.", BLUE, body_size=19)
card(s, 4.85, 1.65, 3.65, 4.45, "Medical-image quality metrics", "Related works use MSE and PSNR to confirm that embedded images remain visually reliable.", CYAN, body_size=19)
card(s, 8.9, 1.65, 3.65, 4.45, "Gap addressed", "This project combines LSB steganography with authentication, browser workflow, and audit logging.", TEAL, body_size=19)
footer(s, 4, TOTAL)

# 5 Motivation
s = prs.slides.add_slide(blank); bg(s); title(s, "Motivation", "Why the system is needed")
bullets(s, [
    "Patient diagnosis notes must remain confidential during image transfer.",
    "Medical images should not be visibly degraded by the security process.",
    "Users need a simple web interface rather than command-line tools.",
    "Operations should be traceable through audit records.",
], .95, 1.8, 11.1, size=25, gap=.78)
footer(s, 5, TOTAL)

# 6 Objectives
s = prs.slides.add_slide(blank); bg(s); title(s, "Objectives", "Project targets")
bullets(s, [
    "Build authenticated access to the medical image system.",
    "Embed patient diagnosis text into PNG/JPG/DICOM-style image workflows.",
    "Recover hidden diagnosis text accurately during decryption.",
    "Preserve visual quality of the steganographic output image.",
    "Store encryption and decryption actions in an audit log.",
], .95, 1.68, 11.4, size=23.5, gap=.68)
footer(s, 6, TOTAL)

# 7 Methodology
s = prs.slides.add_slide(blank); bg(s); title(s, "Methodology", "Least Significant Bit embedding workflow")
card(s, .85, 1.65, 3.65, 3.95, "1. Prepare inputs", "Upload the cover medical image and type the patient diagnosis payload.", BLUE, heading_size=22, body_size=19)
card(s, 4.85, 1.65, 3.65, 3.95, "2. Embed data", "Convert diagnosis text to bits and hide them in the least significant RGB pixel bits.", PURPLE, heading_size=22, body_size=19)
card(s, 8.85, 1.65, 3.65, 3.95, "3. Recover & log", "Upload the stego-image, extract the hidden text, and record the operation in the audit log.", TEAL, heading_size=22, body_size=19)
pill(s, "Cover image", 1.05, 6.05, 1.45, BLUE)
pill(s, "Diagnosis text", 2.75, 6.05, 1.65, CYAN)
pill(s, "LSB algorithm", 4.65, 6.05, 1.65, PURPLE)
pill(s, "Stego PNG", 6.55, 6.05, 1.35, TEAL)
pill(s, "Recovered text", 8.15, 6.05, 1.75, GREEN)
footer(s, 7, TOTAL)

# 8 Results Encryption screenshot
s = prs.slides.add_slide(blank); bg(s); title(s, "Result and Discussion", "Live browser screenshot: Encryption Lab")
enc = ASSET_DIR / "real_encryption_lab_chrome.png"
add_picture_fit(s, enc, .6, 1.55, 8.35, 4.75)
card(s, 9.25, 1.65, 3.15, 1.28, "Real test", "Uploaded demo X-ray and entered patient diagnosis text.", BLUE, heading_size=19, body_size=16)
card(s, 9.25, 3.18, 3.15, 1.28, "Output", "System produced a downloadable secure PNG.", CYAN, heading_size=19, body_size=16)
card(s, 9.25, 4.72, 3.15, 1.28, "Audit", "Encryption event was recorded by the app.", TEAL, heading_size=19, body_size=16)
textbox(s, "Figure 4.1: Live Encryption Lab screenshot showing uploaded X-ray, diagnosis payload, and success state.", .75, 6.55, 11.2, .3, size=15, color=MUTED, align=PP_ALIGN.CENTER)
footer(s, 8, TOTAL)

# 9 Results Decryption screenshot
s = prs.slides.add_slide(blank); bg(s); title(s, "Result and Discussion", "Live browser screenshot: Decryption Interface")
dec = ASSET_DIR / "real_decryption_lab_chrome.png"
add_picture_fit(s, dec, .6, 1.55, 8.35, 4.75)
card(s, 9.25, 1.65, 3.15, 1.28, "Input", "Uploaded the encrypted PNG generated by the browser test.", PURPLE, heading_size=19, body_size=16)
card(s, 9.25, 3.18, 3.15, 1.28, "Recovered", "Diagnosis text was extracted successfully.", GREEN, heading_size=19, body_size=16)
card(s, 9.25, 4.72, 3.15, 1.28, "Verification", "Terminal log shows payload identified and data extracted.", TEAL, heading_size=19, body_size=16)
textbox(s, "Figure 4.2: Live Decryption Lab screenshot displaying successful recovery of hidden diagnosis data.", .75, 6.55, 11.2, .3, size=15, color=MUTED, align=PP_ALIGN.CENTER)
footer(s, 9, TOTAL)

# 10 Contribution
s = prs.slides.add_slide(blank); bg(s); title(s, "Contribution to Knowledge", "What the project adds")
card(s, .85, 1.65, 5.45, 1.55, "Practical secure workflow", "Pairs medical images with diagnosis text without visible external metadata.", BLUE, body_size=18)
card(s, 6.9, 1.65, 5.45, 1.55, "Browser-side prototype", "Demonstrates client-side Canvas processing for steganography in a web app.", CYAN, body_size=18)
card(s, .85, 3.75, 5.45, 1.55, "Authentication + audit", "Adds controlled access and activity records to the steganographic process.", TEAL, body_size=18)
card(s, 6.9, 3.75, 5.45, 1.55, "Defense-ready evidence", "Includes live screenshots, functional testing, and quality discussion.", PURPLE, body_size=18)
footer(s, 10, TOTAL)

# 11 Conclusion & Recommendation
s = prs.slides.add_slide(blank); bg(s); title(s, "Conclusion and Recommendation", "Summary and future work")
card(s, .85, 1.65, 5.6, 4.45, "Conclusion", "The system successfully embeds and retrieves patient diagnosis data from medical images. The live browser test confirmed encryption, decryption, recovery, and audit behavior.", BLUE, body_size=20)
card(s, 6.9, 1.65, 5.6, 4.45, "Recommendations", "• Test with larger DICOM datasets\n• Combine LSB with stronger encryption\n• Add role-based clinical access control\n• Deploy with PostgreSQL and secure cloud storage", TEAL, body_size=18.5)
footer(s, 11, TOTAL)

# 12 References
s = prs.slides.add_slide(blank); bg(s, True); title(s, "References", "Selected sources", True)
refs = [
    "Cheddad, A., Condell, J., Curran, K., & Mc Kevitt, P. (2010). Digital image steganography: survey and analysis.",
    "Morkel, T., Eloff, J. H. P., & Olivier, M. S. (2005). An overview of image steganography.",
    "Provos, N., & Honeyman, P. (2003). Hide and seek: an introduction to steganography.",
    "Sharma, R., & Prabha, S. (2025). Client-side image steganography for clinical data protection.",
    "Kouadri, A., et al. (2025). Quality evaluation of steganographic medical images using MSE and PSNR.",
    "Project materials: Chapter Four.docx, Literature Review PDF, README, and implemented Next.js app.",
]
y = 1.55
for r in refs:
    bullets(s, [r], .85, y, 11.4, size=15.5, color=RGBColor(226,232,240), accent=CYAN, gap=.1)
    y += .72
textbox(s, "Thank you", .9, 6.25, 4.0, .5, size=34, color=WHITE, bold=True, font=FONT_HEAD)
footer(s, 12, TOTAL, True)

prs.save(OUT)
print(OUT)
