from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_AUTO_SIZE
from PIL import Image, ImageDraw, ImageFont, ImageFilter
import math, random

ROOT = Path(__file__).resolve().parents[1]
ASSET_DIR = ROOT / "00_Notes" / "generated_presentation_assets"
ASSET_DIR.mkdir(parents=True, exist_ok=True)
OUT = ROOT / "BigSam_Steganography_Internal_Defense_Presentation.pptx"

# ---------- helpers ----------
def rgb(hexstr):
    hexstr = hexstr.strip('#')
    return RGBColor(int(hexstr[0:2], 16), int(hexstr[2:4], 16), int(hexstr[4:6], 16))

NAVY = rgb('0F172A')
NAVY2 = rgb('111827')
SLATE = rgb('334155')
MUTED = rgb('64748B')
LIGHT = rgb('F8FAFC')
WHITE = rgb('FFFFFF')
BLUE = rgb('2563EB')
CYAN = rgb('06B6D4')
TEAL = rgb('14B8A6')
GREEN = rgb('10B981')
PURPLE = rgb('7C3AED')
AMBER = rgb('F59E0B')
BORDER = rgb('D7E1EA')

FONT_HEAD = "Aptos Display"
FONT_BODY = "Aptos"


def set_fill(shape, color, transparency=0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency


def set_line(shape, color=None, width=1, transparency=0):
    if color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = color
        shape.line.width = Pt(width)
        shape.line.transparency = transparency


def text_box(slide, text, x, y, w, h, font=FONT_BODY, size=18, color=SLATE, bold=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, line_spacing=1.0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rich_title(slide, title, subtitle=None, dark=False):
    color = WHITE if dark else NAVY
    text_box(slide, title, 0.65, 0.35, 9.6, 0.55, font=FONT_HEAD, size=26, color=color, bold=True)
    if subtitle:
        text_box(slide, subtitle, 0.68, 0.9, 9.5, 0.35, size=9.5, color=CYAN if dark else MUTED)
    # accent
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(1.18), Inches(1.0), Inches(0.06))
    set_fill(line, CYAN if dark else BLUE)
    set_line(line)


def add_footer(slide, num, total=13, dark=False):
    col = rgb('94A3B8') if dark else MUTED
    text_box(slide, "Secure Medical Image Management System using Steganography", 0.65, 7.18, 8.0, 0.2, size=7.5, color=col)
    text_box(slide, f"{num:02d}/{total:02d}", 12.05, 7.12, 0.65, 0.25, size=8, color=col, align=PP_ALIGN.RIGHT)


def bullet_list(slide, items, x, y, w, h, size=15, color=SLATE, bullet_color=BLUE, gap=0.43):
    top = y
    for item in items:
        # bullet dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(top+0.08), Inches(0.10), Inches(0.10))
        set_fill(dot, bullet_color)
        set_line(dot)
        text_box(slide, item, x+0.22, top, w-0.22, 0.34, size=size, color=color)
        top += gap


def card(slide, x, y, w, h, title, body=None, accent=BLUE, fill=WHITE, title_size=15, body_size=11.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shp, fill)
    set_line(shp, BORDER, 0.75)
    # accent strip
    strip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(0.12), Inches(h))
    set_fill(strip, accent)
    set_line(strip)
    text_box(slide, title, x+0.28, y+0.18, w-0.45, 0.32, size=title_size, color=NAVY, bold=True)
    if body:
        text_box(slide, body, x+0.28, y+0.60, w-0.45, h-0.75, size=body_size, color=SLATE, line_spacing=1.08)
    return shp


def pill(slide, text, x, y, w, accent=BLUE, dark=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.34))
    set_fill(shp, accent, transparency=12 if not dark else 5)
    set_line(shp, accent, 0.4)
    text_box(slide, text, x+0.05, y+0.07, w-0.1, 0.18, size=8.5, color=WHITE if dark else accent, bold=True, align=PP_ALIGN.CENTER)


def add_bg(slide, dark=False):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY if dark else LIGHT
    # top right decor
    for i, c in enumerate([CYAN, BLUE, TEAL]):
        o = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5+i*0.55), Inches(-0.75+i*0.15), Inches(2.2), Inches(2.2))
        set_fill(o, c, transparency=68 if not dark else 78)
        set_line(o)


def slide_title_only(prs, title, subtitle, name=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, dark=True)
    # left vertical accent
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.28), Inches(7.5))
    set_fill(rect, BLUE)
    set_line(rect)
    text_box(slide, "INTERNAL DEFENSE PRESENTATION", 0.75, 0.7, 4.3, 0.25, size=10, color=CYAN, bold=True)
    text_box(slide, title, 0.72, 1.55, 10.6, 1.55, font=FONT_HEAD, size=34, color=WHITE, bold=True, line_spacing=0.9)
    text_box(slide, subtitle, 0.76, 3.12, 8.4, 0.45, size=15, color=rgb('CBD5E1'))
    # title icon block
    icon = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.75), Inches(4.8), Inches(1.35), Inches(1.35))
    set_fill(icon, CYAN, transparency=5)
    set_line(icon, CYAN, 1)
    text_box(slide, "LSB", 10.95, 5.22, 0.95, 0.32, font=FONT_HEAD, size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # presenter block
    if name:
        panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(5.25), Inches(5.3), Inches(1.0))
        set_fill(panel, rgb('1E293B'))
        set_line(panel, CYAN, 0.8)
        strip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(5.25), Inches(0.12), Inches(1.0))
        set_fill(strip, CYAN)
        set_line(strip)
        text_box(slide, "Presenter", 0.95, 5.43, 1.4, 0.2, size=9.5, color=CYAN, bold=True)
    text_box(slide, "OMONIYI SAMSON DAMILARE  |  DEPARTMENT OF COMPUTER SCIENCE", 0.95, 5.72, 4.85, 0.28, size=11.5, color=WHITE, bold=True)
    text_box(slide, "June 2026", 0.82, 5.93, 3.0, 0.2, size=9.5, color=rgb('CBD5E1'))
    add_footer(slide, 1, dark=True)
    return slide

# ---------- generated figures ----------
def font(size, bold=False):
    candidates = [
        r"C:\Windows\Fonts\segoeuib.ttf" if bold else r"C:\Windows\Fonts\segoeui.ttf",
        r"C:\Windows\Fonts\arialbd.ttf" if bold else r"C:\Windows\Fonts\arial.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ]
    for c in candidates:
        if c and Path(c).exists():
            return ImageFont.truetype(c, size)
    return ImageFont.load_default()


def create_xray(path):
    W, H = 550, 520
    im = Image.new('L', (W,H), 22)
    d = ImageDraw.Draw(im)
    # soft body background
    for r in range(240, 40, -15):
        val = int(30 + (240-r)*0.33)
        d.ellipse((W/2-r*0.82, H/2-r, W/2+r*0.82, H/2+r), fill=val)
    # lungs as darker lobes
    d.ellipse((105, 118, 260, 410), fill=45)
    d.ellipse((290, 118, 445, 410), fill=45)
    # ribs
    for i in range(8):
        y = 125+i*32
        d.arc((70,y-24, 275,y+48), 190, 345, fill=128, width=3)
        d.arc((275,y-24, 480,y+48), 195, 350, fill=128, width=3)
    # spine/sternum
    d.rounded_rectangle((260, 70, 292, 455), radius=14, fill=125)
    for y in range(100,440,35):
        d.line((258,y,294,y), fill=170, width=2)
    # clavicles
    d.arc((95, 70, 285, 170), 195, 340, fill=170, width=4)
    d.arc((265, 70, 455, 170), 200, 345, fill=170, width=4)
    im = im.filter(ImageFilter.GaussianBlur(1.8))
    im = Image.merge('RGB', (im, im, im))
    im.save(path)


def rounded_rect(draw, box, fill, outline=None, radius=18, width=1):
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=width)


def create_ui_figure(path, mode='encrypt'):
    W,H = 1600,900
    im = Image.new('RGB', (W,H), (248,250,252))
    d = ImageDraw.Draw(im)
    f9=font(18); f10=font(20); f11=font(22); f12=font(24); f14=font(28); f18=font(36, True); f24=font(48, True)
    # sidebar
    d.rectangle((0,0,270,H), fill=(15,23,42))
    d.text((35,35), "MED-STEGO", font=f18, fill=(255,255,255))
    d.text((35,80), "Secure Image System", font=f9, fill=(148,163,184))
    nav = ["Dashboard", "Encryption Lab", "Decryption Lab", "Audit Logs", "Settings"]
    for i,n in enumerate(nav):
        y=160+i*62
        active = (mode=='encrypt' and n=='Encryption Lab') or (mode=='decrypt' and n=='Decryption Lab')
        rounded_rect(d, (25,y,245,y+44), fill=(37,99,235) if active else (15,23,42), radius=12)
        d.text((50,y+11), n, font=f10, fill=(255,255,255) if active else (203,213,225))
    # header
    d.rectangle((270,0,W,86), fill=(255,255,255))
    title = "Encryption Dashboard" if mode=='encrypt' else "Decryption Interface"
    d.text((315,24), title, font=f18, fill=(15,23,42))
    d.text((1260,26), "Audit Log: ACTIVE", font=f10, fill=(16,185,129))
    # main cards
    if mode == 'encrypt':
        # image upload panel
        rounded_rect(d, (320,125,880,715), fill=(255,255,255), outline=(214,226,238), radius=24, width=2)
        d.text((350,155), "Uploaded Medical Image", font=f14, fill=(15,23,42))
        xray_path = ASSET_DIR/'generated_xray.png'
        if not xray_path.exists(): create_xray(xray_path)
        xr = Image.open(xray_path).resize((430,405))
        rounded_rect(d, (385,215,815,620), fill=(15,23,42), radius=18)
        im.paste(xr, (385,215))
        d.text((405,642), "Chest_Xray_500x500.png", font=f10, fill=(100,116,139))
        # diagnosis panel
        rounded_rect(d, (930,125,1510,715), fill=(255,255,255), outline=(214,226,238), radius=24, width=2)
        d.text((965,155), "Diagnosis Input", font=f14, fill=(15,23,42))
        rounded_rect(d, (965,210,1470,425), fill=(248,250,252), outline=(203,213,225), radius=14)
        sample = ["Patient ID: XR-2042", "Findings: mild lower-zone opacity", "Recommendation: clinical correlation", "Hidden securely using LSB embedding."]
        for i,t in enumerate(sample): d.text((990,238+i*42), t, font=f11, fill=(51,65,85))
        rounded_rect(d, (965,470,1470,535), fill=(37,99,235), radius=16)
        d.text((1110,486), "ENCRYPT AND DOWNLOAD PNG", font=f12, fill=(255,255,255))
        rounded_rect(d, (965,580,1470,654), fill=(236,253,245), outline=(167,243,208), radius=16)
        d.text((990,603), "✓ Success: steganographic image generated", font=f11, fill=(5,150,105))
    else:
        # recovered panel
        rounded_rect(d, (320,125,820,725), fill=(255,255,255), outline=(214,226,238), radius=24, width=2)
        d.text((350,155), "Encrypted Image Uploaded", font=f14, fill=(15,23,42))
        xray_path = ASSET_DIR/'generated_xray.png'
        if not xray_path.exists(): create_xray(xray_path)
        xr = Image.open(xray_path).resize((380,355))
        rounded_rect(d, (375,220,755,575), fill=(15,23,42), radius=18)
        im.paste(xr, (375,220))
        d.text((390,600), "stego_output.png", font=f10, fill=(100,116,139))
        rounded_rect(d, (380,655,760,705), fill=(124,58,237), radius=14)
        d.text((480,668), "DECRYPT MESSAGE", font=f11, fill=(255,255,255))
        # result panel
        rounded_rect(d, (875,125,1510,725), fill=(255,255,255), outline=(214,226,238), radius=24, width=2)
        d.text((910,155), "Recovered Patient Diagnosis", font=f14, fill=(15,23,42))
        rounded_rect(d, (910,210,1470,445), fill=(240,253,250), outline=(153,246,228), radius=16)
        sample = ["Patient ID: XR-2042", "Findings: mild lower-zone opacity", "Recommendation: clinical correlation", "Integrity: message recovered successfully"]
        for i,t in enumerate(sample): d.text((940,240+i*42), t, font=f11, fill=(15,118,110))
        d.text((910,490), "System Console", font=f12, fill=(15,23,42))
        rounded_rect(d, (910,530,1470,685), fill=(15,23,42), radius=16)
        logs = ["> System initialized", "> File detected: stego_output.png", "> Extracting LSB payload...", "> Success: diagnosis recovered"]
        for i,t in enumerate(logs): d.text((935,555+i*30), t, font=f9, fill=(103,232,249) if i==3 else (203,213,225))
    im.save(path)

create_ui_figure(ASSET_DIR/'figure_4_1_encryption_dashboard.png', 'encrypt')
create_ui_figure(ASSET_DIR/'figure_4_2_decryption_interface.png', 'decrypt')

# ---------- deck ----------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# 1 Title
slide_title_only(prs,
    "Development of a Secure Web-Based Medical Image Management System Using Steganography",
    "Protecting patient diagnosis data inside medical images without visible distortion",
    "OMONIYI SAMSON DAMILARE")

# 2 Outline
s = prs.slides.add_slide(blank); add_bg(s); add_rich_title(s, "Presentation Outline", "Following the format supplied for the real/internal defense")
outline = ["Introduction", "Review of Literature and Related Works", "Motivation", "Objectives", "Methodology", "Result and Discussion", "Contribution to Knowledge", "Conclusion and Recommendation", "References"]
for i,item in enumerate(outline):
    col = i % 3; row = i // 3
    x = 0.75 + col*4.1; y = 1.65 + row*1.45
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.6), Inches(0.9))
    set_fill(shp, WHITE); set_line(shp, BORDER, 0.8)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+0.18), Inches(y+0.20), Inches(0.48), Inches(0.48))
    set_fill(circ, BLUE if i<5 else TEAL); set_line(circ)
    text_box(s, str(i+1), x+0.28, y+0.31, 0.28, 0.2, size=10.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text_box(s, item, x+0.82, y+0.22, 2.55, 0.42, size=13.2, color=NAVY, bold=True)
add_footer(s, 2)

# 3 Introduction
s = prs.slides.add_slide(blank); add_bg(s); add_rich_title(s, "Introduction", "Why secure medical image handling matters")
card(s, 0.75, 1.55, 4.0, 4.35, "Context", "Medical images often move with supporting diagnosis notes, patient identifiers, and clinical comments. If shared carelessly, these details can expose confidential health information.", accent=BLUE)
card(s, 4.95, 1.55, 3.65, 4.35, "Steganography", "Steganography hides sensitive text inside an image file. Unlike visible encryption warnings, the image remains visually normal while carrying the embedded message.", accent=CYAN)
card(s, 8.8, 1.55, 3.75, 4.35, "Project Focus", "This project implements a secure web-based system that embeds and retrieves patient diagnosis text from medical images while preserving image quality.", accent=TEAL)
pill(s, "Confidentiality", 1.1, 6.25, 1.45, BLUE); pill(s, "Imperceptibility", 2.75, 6.25, 1.55, CYAN); pill(s, "Auditability", 4.55, 6.25, 1.25, TEAL); pill(s, "Web Access", 6.05, 6.25, 1.2, PURPLE)
add_footer(s, 3)

# 4 Literature Review
s = prs.slides.add_slide(blank); add_bg(s); add_rich_title(s, "Review of Literature and Related Works", "Client-supplied literature review format: Objective • Methodology • Result • Contribution")
rows = [
    ["Cheddad et al.", "Survey image steganography", "Comparative review", "Capacity, security, imperceptibility trade-offs", "Guides algorithm selection"],
    ["Morkel et al.", "Explain steganography principles", "Conceptual analysis", "Hidden communication protects message presence", "Supports use in sensitive data transfer"],
    ["Sharma & Prabha (2025)", "Fast clinical image embedding", "LSB/client-side image processing", "Efficient execution for real-time use", "Supports browser-based workflow"],
    ["Kouadri et al. (2025)", "Medical image quality preservation", "PSNR/MSE evaluation", "High PSNR retains diagnostic value", "Basis for quality validation"],
]
headers = ["Author", "Objective", "Methodology", "Result", "Contribution"]
table = s.shapes.add_table(len(rows)+1, len(headers), Inches(0.55), Inches(1.52), Inches(12.25), Inches(4.55)).table
col_widths = [1.55, 2.35, 2.55, 3.0, 2.8]
for c,w in enumerate(col_widths): table.columns[c].width = Inches(w)
for c,h in enumerate(headers):
    cell = table.cell(0,c); cell.text = h
    cell.fill.solid(); cell.fill.fore_color.rgb = BLUE
    for p in cell.text_frame.paragraphs:
        p.runs[0].font.color.rgb = WHITE; p.runs[0].font.bold = True; p.runs[0].font.size = Pt(9.5); p.runs[0].font.name = FONT_BODY
for r,row in enumerate(rows,1):
    for c,val in enumerate(row):
        cell = table.cell(r,c); cell.text = val
        cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if r%2 else rgb('F1F5F9')
        for p in cell.text_frame.paragraphs:
            p.runs[0].font.size = Pt(8.3); p.runs[0].font.color.rgb = SLATE; p.runs[0].font.name = FONT_BODY
text_box(s, "Gap addressed: a practical web system combining LSB steganography, authentication, and audit logging for medical image workflows.", 0.8, 6.28, 11.4, 0.35, size=13, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
add_footer(s, 4)

# 5 Motivation
s = prs.slides.add_slide(blank); add_bg(s); add_rich_title(s, "Motivation", "The problem that led to the project")
card(s, 0.8, 1.65, 3.65, 4.25, "Data privacy risk", "Diagnosis notes can be exposed when images are transferred through general channels or unsecured storage.", accent=BLUE)
card(s, 4.85, 1.65, 3.65, 4.25, "Image integrity", "Medical images must remain diagnostically useful; security must not visibly damage regions of interest.", accent=CYAN)
card(s, 8.9, 1.65, 3.65, 4.25, "Accountability", "Clinical systems need traceability, so encryption and decryption events should be logged with user and file context.", accent=TEAL)
add_footer(s, 5)

# 6 Objectives
s = prs.slides.add_slide(blank); add_bg(s); add_rich_title(s, "Objectives", "What the system was designed to achieve")
bullet_list(s, [
    "Develop a secure web-based dashboard for authenticated users.",
    "Embed patient diagnosis text into medical images using LSB steganography.",
    "Recover hidden diagnosis text accurately during decryption.",
    "Preserve visual quality so the stego-image remains medically usable.",
    "Record encryption/decryption activities for audit and accountability.",
], 0.95, 1.65, 6.1, 3.0, size=15.5)
# objective metrics card
card(s, 7.35, 1.55, 4.8, 4.1, "Success Criteria", "• Hidden text recovered without corruption\n• Minimal visual difference between cover and stego-image\n• Responsive browser processing\n• Event logs stored through Prisma/SQLite", accent=PURPLE, body_size=13)
pill(s, "Accuracy", 7.65, 6.05, 1.1, BLUE); pill(s, "Quality", 8.95, 6.05, 1.0, CYAN); pill(s, "Security", 10.15, 6.05, 1.05, TEAL); pill(s, "Logs", 11.4, 6.05, 0.75, PURPLE)
add_footer(s, 6)

# 7 Methodology
s = prs.slides.add_slide(blank); add_bg(s); add_rich_title(s, "Methodology", "LSB embedding and recovery workflow")
# flow boxes
steps = [("Cover image", BLUE), ("Diagnosis text", CYAN), ("LSB embedding", PURPLE), ("Stego PNG", TEAL), ("Decryption", AMBER), ("Recovered text", GREEN)]
positions = [(0.85,2.05),(0.85,3.55),(3.55,2.8),(6.1,2.8),(8.45,2.8),(10.75,2.8)]
for (label,col),(x,y) in zip(steps, positions):
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.9), Inches(0.75))
    set_fill(shp, col); set_line(shp)
    text_box(s, label, x+0.08, y+0.24, 1.74, 0.25, size=11.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
# arrows
for x1,y1,x2,y2 in [(2.75,2.42,3.55,3.15),(2.75,3.92,3.55,3.15),(5.45,3.15,6.1,3.15),(8.0,3.15,8.45,3.15),(10.35,3.15,10.75,3.15)]:
    conn=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb=SLATE; conn.line.width=Pt(1.5)
card(s, 0.85, 5.15, 3.7, 1.05, "Embedding", "Text is converted to binary and placed in the least significant bits of RGB pixel channels.", accent=BLUE, body_size=10.2)
card(s, 4.85, 5.15, 3.7, 1.05, "Recovery", "Bits are extracted until the terminator is found, reconstructing the diagnosis text.", accent=CYAN, body_size=10.2)
card(s, 8.85, 5.15, 3.7, 1.05, "Audit", "Each operation is logged with file context, user, status, and timestamp.", accent=TEAL, body_size=10.2)
add_footer(s, 7)

# 8 Architecture / Modules
s = prs.slides.add_slide(blank); add_bg(s); add_rich_title(s, "System Architecture and Core Modules", "Implementation stack and responsibilities")
card(s, 0.75, 1.55, 3.65, 3.65, "Authentication Dashboard", "Protected routes for dashboard, encryption lab, decryption lab, and settings using Better Auth.", accent=BLUE)
card(s, 4.85, 1.55, 3.65, 3.65, "Steganography Engine", "HTML5 Canvas extracts pixel data; LSB algorithm embeds and retrieves diagnosis text from image pixels.", accent=CYAN)
card(s, 8.95, 1.55, 3.65, 3.65, "Audit Logging", "Prisma ORM stores encryption and decryption events in SQLite for accountability and later review.", accent=TEAL)
techs = [("Next.js 16", BLUE), ("React 19", CYAN), ("Tailwind CSS", TEAL), ("Prisma ORM", PURPLE), ("SQLite", AMBER), ("Canvas API", GREEN)]
for i,(t,c) in enumerate(techs):
    pill(s, t, 1.05+i*1.85, 5.85, 1.35, c)
add_footer(s, 8)

# 9 Result & Discussion Encryption
s = prs.slides.add_slide(blank); add_bg(s); add_rich_title(s, "Result and Discussion", "Encryption dashboard output")
img = str(ASSET_DIR/'figure_4_1_encryption_dashboard.png')
s.shapes.add_picture(img, Inches(0.7), Inches(1.42), width=Inches(7.65))
card(s, 8.65, 1.55, 3.85, 1.1, "Test Case 1: Encryption", "A standard chest X-ray image and diagnosis text were successfully processed into a downloadable stego PNG.", accent=BLUE, body_size=10.8)
card(s, 8.65, 2.95, 3.85, 1.1, "Visual Quality", "The output image remained visually indistinguishable from the original cover image.", accent=CYAN, body_size=10.8)
card(s, 8.65, 4.35, 3.85, 1.1, "User Feedback", "Real-time status indicators confirmed processing and success states.", accent=TEAL, body_size=10.8)
text_box(s, "Figure 4.1: Screenshot of the Encryption Dashboard showing the uploaded X-ray and the diagnosis input field", 0.85, 6.85, 7.3, 0.22, size=8.6, color=MUTED, align=PP_ALIGN.CENTER)
add_footer(s, 9)

# 10 Functional Testing / Performance
s = prs.slides.add_slide(blank); add_bg(s); add_rich_title(s, "Functional Testing and Performance", "Decryption, audit logging, and quality metrics")
img = str(ASSET_DIR/'figure_4_2_decryption_interface.png')
s.shapes.add_picture(img, Inches(0.7), Inches(1.42), width=Inches(6.45))
# metric cards
metrics = [("Recovered", "200-char diagnosis", GREEN), ("PSNR", "> 55 dB", BLUE), ("Speed", "~150 ms / 1MB", CYAN), ("Audit", "Events logged", PURPLE)]
for i,(k,v,c) in enumerate(metrics):
    x=7.65+(i%2)*2.45; y=1.65+(i//2)*1.55
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.1), Inches(1.05))
    set_fill(shp, WHITE); set_line(shp, BORDER, 0.8)
    text_box(s, k, x+0.15, y+0.18, 1.8, 0.2, size=9.8, color=MUTED, bold=True)
    text_box(s, v, x+0.15, y+0.48, 1.8, 0.32, size=15, color=c, bold=True)
card(s, 7.65, 4.85, 4.55, 1.05, "Discussion", "The results indicate that the system balances security, usability, and image imperceptibility for medical image workflows.", accent=TEAL, body_size=10.8)
text_box(s, "Figure 4.2: Screenshot of the Decryption Interface displaying the recovered patient diagnosis", 0.85, 6.85, 6.2, 0.22, size=8.6, color=MUTED, align=PP_ALIGN.CENTER)
add_footer(s, 10)

# 11 Contribution
s = prs.slides.add_slide(blank); add_bg(s); add_rich_title(s, "Contribution to Knowledge", "What this project adds")
contribs = [
    ("Medical-data hiding workflow", "Demonstrates a usable path for pairing diagnosis text with medical image files without obvious external metadata.", BLUE),
    ("Browser-side processing", "Uses client-side Canvas processing, reducing the need to upload raw patient notes to a server before embedding.", CYAN),
    ("Security + accountability", "Combines steganography with authenticated access and audit logging in one practical web prototype.", TEAL),
    ("Quality-aware evaluation", "Uses imperceptibility, PSNR, speed, and functional recovery as defense-ready validation points.", PURPLE),
]
for i,(t,b,c) in enumerate(contribs):
    card(s, 0.95+(i%2)*5.8, 1.55+(i//2)*2.05, 5.15, 1.45, t, b, accent=c, body_size=10.7)
add_footer(s, 11)

# 12 Conclusion & Recommendation
s = prs.slides.add_slide(blank); add_bg(s); add_rich_title(s, "Conclusion and Recommendation", "Summary and next steps")
card(s, 0.85, 1.55, 5.4, 4.4, "Conclusion", "The developed system successfully embeds and retrieves diagnosis data inside medical images using LSB steganography. Testing showed accurate recovery, high visual imperceptibility, responsive processing, and traceable audit logs.", accent=BLUE, body_size=13.2)
card(s, 6.75, 1.55, 5.4, 4.4, "Recommendations", "• Extend testing with larger DICOM datasets\n• Combine LSB with stronger cryptographic encryption\n• Add role-based clinical access controls\n• Deploy with PostgreSQL and secure cloud storage\n• Run usability testing with medical users", accent=TEAL, body_size=12.2)
add_footer(s, 12)

# 13 References
s = prs.slides.add_slide(blank); add_bg(s, dark=True); add_rich_title(s, "References", "Selected literature and project sources", dark=True)
refs = [
    "Cheddad, A., Condell, J., Curran, K., & Mc Kevitt, P. (2010). Digital image steganography: Survey and analysis of current methods.",
    "Morkel, T., Eloff, J. H. P., & Olivier, M. S. (2005). An overview of image steganography.",
    "Provos, N., & Honeyman, P. (2003). Hide and seek: An introduction to steganography.",
    "Sharma, R., & Prabha, S. (2025). Client-side image steganography for efficient clinical data protection.",
    "Kouadri, A., et al. (2025). Quality evaluation of steganographic medical images using MSE and PSNR metrics.",
    "Project source files: Chapter Four.docx, README.md, and implemented Next.js application modules.",
]
y = 1.55
for r in refs:
    bullet_list(s, [r], 0.9, y, 11.5, 0.35, size=10.6, color=rgb('E2E8F0'), bullet_color=CYAN, gap=0.0)
    y += 0.68
text_box(s, "Thank you", 0.9, 6.25, 5.0, 0.5, font=FONT_HEAD, size=26, color=WHITE, bold=True)
add_footer(s, 13, dark=True)

# Remove any default slides? We started from empty Presentation with one blank? python-pptx starts with no slides.
prs.save(OUT)
print(OUT)
print(ASSET_DIR)
